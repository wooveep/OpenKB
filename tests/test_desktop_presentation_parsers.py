"""Focused PPTX import coverage for the Desktop Document IR boundary."""

from __future__ import annotations

import json
import sqlite3
from io import BytesIO
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from PIL import Image as PillowImage
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches

from openkb.desktop_import import DesktopTextImportService
from openkb.desktop_import_artifacts import DesktopImportError
from openkb.desktop_raw_assets import DesktopRawAssetService
from openkb.desktop_workspace import DesktopKnowledgeBaseRuntime


def test_pptx_import_retains_slide_order_shapes_notes_tables_and_source_images(tmp_path):
    kb_dir = tmp_path / "desktop-kb"
    source = tmp_path / "briefing.pptx"
    image_bytes = _presentation_with_slides_and_image()
    source.write_bytes(image_bytes)
    DesktopKnowledgeBaseRuntime().create(kb_dir)

    imported = DesktopTextImportService(kb_dir).import_text(source)

    assert imported.document.source_format == "pptx"
    raw_path = kb_dir / "raw" / f"{imported.document.raw_asset_sha256}.pptx"
    assert raw_path.read_bytes() == image_bytes
    with sqlite3.connect(kb_dir / ".openkb" / "state.sqlite3") as connection:
        rows = connection.execute(
            "SELECT kind, text, locator_json FROM document_ir_blocks ORDER BY ordinal"
        ).fetchall()
        image_row = connection.execute(
            "SELECT storage_path, locator_json FROM source_images"
        ).fetchone()

    slide_one_blocks = [
        (kind, text, json.loads(locator))
        for kind, text, locator in rows
        if json.loads(locator).get("slide_index") == 1
    ]
    texts = [text for _, text, _ in slide_one_blocks]
    assert texts.index("Overview") < texts.index("Right-side text")
    assert any(kind == "list" and "- First bullet" in text for kind, text, _ in slide_one_blocks)
    assert any(
        kind == "paragraph" and text == "Lead-in text\n  Indented non-list"
        for kind, text, _ in slide_one_blocks
    )
    assert any(kind == "table" and "| Name | Value |" in text for kind, text, _ in slide_one_blocks)
    notes = [
        locator for _, _, locator in slide_one_blocks if locator.get("source") == "speaker_notes"
    ]
    assert notes == [{"slide": 1, "slide_index": 1, "source": "speaker_notes"}]
    assert image_row is not None
    image_path, image_locator_json = image_row
    assert json.loads(image_locator_json)["slide_index"] == 1
    assert json.loads(image_locator_json)["shape_index"]
    assert (kb_dir / image_path).read_bytes().startswith(b"\x89PNG")

    reader = DesktopRawAssetService(kb_dir).read_document(imported.document.document_id)
    assert "# Slide 1" in reader.content
    assert "# Slide 2" in reader.content
    assert "Presenter note" in reader.content
    assert len(reader.source_images) == 1


def test_pptx_with_a_malformed_slide_has_a_stable_import_failure(tmp_path):
    kb_dir = tmp_path / "desktop-kb"
    source = tmp_path / "broken-slide.pptx"
    source.write_bytes(_presentation_with_malformed_slide())
    DesktopKnowledgeBaseRuntime().create(kb_dir)

    with pytest.raises(DesktopImportError) as error:
        DesktopTextImportService(kb_dir).import_text(source)

    assert error.value.code == "invalid_pptx_document"


def _presentation_with_slides_and_image() -> bytes:
    presentation = Presentation()
    _add_inherited_bullet_style(presentation)
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    right = slide.shapes.add_textbox(Inches(5), Inches(1), Inches(2), Inches(1))
    right.text_frame.text = "Right-side text"
    overview = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    overview.text_frame.text = "Overview"

    bullets = slide.placeholders[1]
    bullets_frame = bullets.text_frame
    bullets_frame.paragraphs[0].text = "First bullet"
    nested = bullets_frame.add_paragraph()
    nested.text = "Nested bullet"
    nested.level = 1

    non_list = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    non_list_frame = non_list.text_frame
    non_list_frame.paragraphs[0].text = "Lead-in text"
    indented = non_list_frame.add_paragraph()
    indented.text = "Indented non-list"
    indented.level = 1

    table = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "OpenKB"
    table.cell(1, 1).text = "Desktop"

    image = BytesIO()
    PillowImage.new("RGB", (1, 1), "blue").save(image, format="PNG")
    slide.shapes.add_picture(BytesIO(image.getvalue()), Inches(1), Inches(4))
    slide.notes_slide.notes_text_frame.text = "Presenter note"

    second_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    second_slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(4), Inches(1)
    ).text_frame.text = "Second slide text"

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _add_inherited_bullet_style(presentation: Presentation) -> None:
    content_placeholder = presentation.slide_layouts[1].placeholders[1]
    text_body = content_placeholder._element.find(qn("p:txBody"))
    assert text_body is not None
    list_style = text_body.find(qn("a:lstStyle"))
    assert list_style is not None
    level_one = OxmlElement("a:lvl1pPr")
    bullet = OxmlElement("a:buChar")
    bullet.set("char", "•")
    level_one.append(bullet)
    list_style.append(level_one)


def _presentation_with_malformed_slide() -> bytes:
    valid = _presentation_with_slides_and_image()
    malformed = BytesIO()
    with (
        ZipFile(BytesIO(valid)) as source,
        ZipFile(malformed, "w", compression=ZIP_DEFLATED) as destination,
    ):
        for item in source.infolist():
            content = source.read(item.filename)
            if item.filename == "ppt/slides/slide1.xml":
                content = b"<p:sld>"
            destination.writestr(item, content)
    return malformed.getvalue()
