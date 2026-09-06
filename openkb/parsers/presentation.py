"""Structure-preserving PPTX adapter for Desktop document import."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any
from uuid import uuid4
from zipfile import BadZipFile

from lxml.etree import XMLSyntaxError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.exc import PackageNotFoundError
from pptx.oxml.ns import qn

from openkb.importing.artifacts import (
    DesktopImportError,
    DocumentIRBlock,
    ParsedDocument,
    SourceImage,
    source_image_from_content,
)


def parse_presentation_document(source: Path, raw_bytes: bytes) -> ParsedDocument:
    """Parse a PPTX Raw Asset into citable slide-level Document IR."""
    try:
        presentation = Presentation(BytesIO(raw_bytes))
        blocks: list[DocumentIRBlock] = []
        images: list[SourceImage] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            _append_slide(slide, slide_index=slide_index, blocks=blocks, images=images)
    except (
        BadZipFile,
        KeyError,
        OSError,
        PackageNotFoundError,
        ValueError,
        XMLSyntaxError,
    ) as error:
        raise DesktopImportError(
            "invalid_pptx_document", f"PPTX source cannot be read: {source.name}"
        ) from error
    if not blocks:
        raise DesktopImportError(
            "empty_text_document", f"PPTX source did not contain usable slides: {source.name}"
        )
    return ParsedDocument(tuple(blocks), tuple(images))


def _append_slide(
    slide: Any,
    *,
    slide_index: int,
    blocks: list[DocumentIRBlock],
    images: list[SourceImage],
) -> None:
    slide_name = f"Slide {slide_index}"
    heading_path = (slide_name,)
    blocks.append(
        DocumentIRBlock(
            block_id=_block_id(),
            ordinal=len(blocks),
            kind="heading",
            text=slide_name,
            heading_path=heading_path,
            line_start=slide_index,
            line_end=slide_index,
            locator={"slide": slide_index, "slide_index": slide_index, "heading_level": 1},
        )
    )
    for shape, shape_path in _ordered_shapes(slide.shapes):
        locator = _shape_locator(slide_index, shape, shape_path)
        if bool(getattr(shape, "has_table", False)):
            _append_table_shape(shape, heading_path, locator, slide_index, blocks)
        elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            _append_picture_shape(
                shape, shape_path, heading_path, locator, slide_index, blocks, images
            )
        elif bool(getattr(shape, "has_text_frame", False)):
            _append_text_shape(shape, heading_path, locator, slide_index, blocks)
    _append_speaker_notes(slide, heading_path, slide_index, blocks)


def _ordered_shapes(
    shapes: Any, path: tuple[int, ...] = ()
) -> tuple[tuple[Any, tuple[int, ...]], ...]:
    ordered = sorted(
        enumerate(shapes, start=1),
        key=lambda item: (
            int(getattr(item[1], "top", 0) or 0),
            int(getattr(item[1], "left", 0) or 0),
            item[0],
        ),
    )
    flattened: list[tuple[Any, tuple[int, ...]]] = []
    for shape_index, shape in ordered:
        shape_path = path + (shape_index,)
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            flattened.extend(_ordered_shapes(shape.shapes, shape_path))
        else:
            flattened.append((shape, shape_path))
    return tuple(flattened)


def _shape_locator(slide_index: int, shape: Any, shape_path: tuple[int, ...]) -> dict[str, object]:
    locator: dict[str, object] = {
        "slide": slide_index,
        "slide_index": slide_index,
        "shape_index": shape_path[0],
        "shape_path": list(shape_path),
        "geometry": {
            "left": int(getattr(shape, "left", 0) or 0),
            "top": int(getattr(shape, "top", 0) or 0),
            "width": int(getattr(shape, "width", 0) or 0),
            "height": int(getattr(shape, "height", 0) or 0),
        },
    }
    shape_name = str(getattr(shape, "name", "")).strip()
    if shape_name:
        locator["shape_name"] = shape_name
    return locator


def _append_text_shape(
    shape: Any,
    heading_path: tuple[str, ...],
    locator: dict[str, object],
    slide_index: int,
    blocks: list[DocumentIRBlock],
) -> None:
    current_kind: str | None = None
    current_start = 0
    current_values: list[str] = []
    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
        value = paragraph.text.strip()
        if not value:
            continue
        level = int(paragraph.level)
        is_list_item = _paragraph_is_list_item(shape, paragraph)
        kind = "list" if is_list_item else "paragraph"
        rendered = f"{'  ' * level}- {value}" if is_list_item else f"{'  ' * level}{value}"
        if current_kind is not None and current_kind != kind:
            _append_text_block(
                current_kind,
                current_values,
                current_start,
                paragraph_index - 1,
                heading_path,
                locator,
                slide_index,
                blocks,
            )
            current_values = []
        if not current_values:
            current_kind = kind
            current_start = paragraph_index
        current_values.append(rendered)
    if current_kind is not None:
        _append_text_block(
            current_kind,
            current_values,
            current_start,
            current_start + len(current_values) - 1,
            heading_path,
            locator,
            slide_index,
            blocks,
        )


def _append_text_block(
    kind: str,
    values: list[str],
    paragraph_start: int,
    paragraph_end: int,
    heading_path: tuple[str, ...],
    locator: dict[str, object],
    slide_index: int,
    blocks: list[DocumentIRBlock],
) -> None:
    blocks.append(
        DocumentIRBlock(
            block_id=_block_id(),
            ordinal=len(blocks),
            kind=kind,
            text="\n".join(values),
            heading_path=heading_path,
            line_start=slide_index,
            line_end=slide_index,
            locator={
                **locator,
                "paragraph_start": paragraph_start,
                "paragraph_end": paragraph_end,
            },
        )
    )


def _paragraph_is_list_item(shape: Any, paragraph: Any) -> bool:
    state = _bullet_state(getattr(paragraph._p, "pPr", None))
    if state is not None:
        return state
    for list_style in _inherited_list_styles(shape):
        state = _list_style_bullet_state(list_style, int(paragraph.level))
        if state is not None:
            return state
    return False


def _bullet_state(paragraph_properties: Any) -> bool | None:
    if paragraph_properties is None:
        return None
    child_names = {str(child.tag).rsplit("}", maxsplit=1)[-1] for child in paragraph_properties}
    if "buNone" in child_names:
        return False
    if {"buAutoNum", "buBlip", "buChar"} & child_names:
        return True
    return None


def _inherited_list_styles(shape: Any) -> tuple[Any, ...]:
    list_styles: list[Any] = [_text_list_style(shape)]
    if not bool(getattr(shape, "is_placeholder", False)):
        return tuple(style for style in list_styles if style is not None)
    placeholder_index = int(shape.placeholder_format.idx)
    placeholder_type = shape.placeholder_format.type
    layout = shape.part.slide_layout
    layout_placeholder = _matching_placeholder(
        layout.placeholders, placeholder_index, placeholder_type
    )
    if layout_placeholder is not None:
        list_styles.append(_text_list_style(layout_placeholder))
    master = layout.slide_master
    master_placeholder = _matching_placeholder(
        master.placeholders, placeholder_index, placeholder_type
    )
    if master_placeholder is not None:
        list_styles.append(_text_list_style(master_placeholder))
    master_style = _master_text_style(master, placeholder_type)
    if master_style is not None:
        list_styles.append(master_style)
    return tuple(style for style in list_styles if style is not None)


def _text_list_style(shape: Any) -> Any | None:
    text_body = shape._element.find(qn("p:txBody"))
    return None if text_body is None else text_body.find(qn("a:lstStyle"))


def _matching_placeholder(
    placeholders: Any, placeholder_index: int, placeholder_type: Any
) -> Any | None:
    matches = tuple(placeholders)
    for placeholder in matches:
        if int(placeholder.placeholder_format.idx) == placeholder_index:
            return placeholder
    for placeholder in matches:
        if placeholder.placeholder_format.type == placeholder_type:
            return placeholder
    return None


def _master_text_style(master: Any, placeholder_type: Any) -> Any | None:
    text_styles = master._element.find(qn("p:txStyles"))
    if text_styles is None:
        return None
    if placeholder_type in {
        PP_PLACEHOLDER.TITLE,
        PP_PLACEHOLDER.CENTER_TITLE,
        PP_PLACEHOLDER.SUBTITLE,
        PP_PLACEHOLDER.VERTICAL_TITLE,
    }:
        style_name = "titleStyle"
    elif placeholder_type in {
        PP_PLACEHOLDER.BODY,
        PP_PLACEHOLDER.OBJECT,
        PP_PLACEHOLDER.VERTICAL_BODY,
        PP_PLACEHOLDER.VERTICAL_OBJECT,
    }:
        style_name = "bodyStyle"
    else:
        style_name = "otherStyle"
    return text_styles.find(qn(f"p:{style_name}"))


def _list_style_bullet_state(list_style: Any, level: int) -> bool | None:
    if not 0 <= level <= 8:
        return None
    paragraph_properties = list_style.find(qn(f"a:lvl{level + 1}pPr"))
    return _bullet_state(paragraph_properties)


def _append_table_shape(
    shape: Any,
    heading_path: tuple[str, ...],
    locator: dict[str, object],
    slide_index: int,
    blocks: list[DocumentIRBlock],
) -> None:
    values = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
    if not values or not any(value for row in values for value in row):
        return
    table_locator = {
        **locator,
        "row_count": len(values),
        "column_count": max((len(row) for row in values), default=0),
        "cell_range": f"R1C1:R{len(values)}C{max((len(row) for row in values), default=0)}",
    }
    blocks.append(
        DocumentIRBlock(
            block_id=_block_id(),
            ordinal=len(blocks),
            kind="table",
            text=_markdown_table(values),
            heading_path=heading_path,
            line_start=slide_index,
            line_end=slide_index,
            locator=table_locator,
        )
    )


def _append_picture_shape(
    shape: Any,
    shape_path: tuple[int, ...],
    heading_path: tuple[str, ...],
    locator: dict[str, object],
    slide_index: int,
    blocks: list[DocumentIRBlock],
    images: list[SourceImage],
) -> None:
    try:
        picture = shape.image
        content = picture.blob
    except (AttributeError, KeyError, OSError, ValueError):
        return
    if not content:
        return
    extension = str(getattr(picture, "ext", "") or "bin").lstrip(".")
    shape_path_text = ".".join(str(value) for value in shape_path)
    filename = f"slide-{slide_index}-shape-{shape_path_text}.{extension}"
    alt_text = str(getattr(shape, "name", "")).strip() or None
    source_image = source_image_from_content(
        content=content,
        filename=filename,
        alt_text=alt_text,
        ordinal=len(images),
        locator=locator,
    )
    images.append(source_image)
    blocks.append(
        DocumentIRBlock(
            block_id=_block_id(),
            ordinal=len(blocks),
            kind="figure",
            text=source_image.alt_text or filename,
            heading_path=heading_path,
            line_start=slide_index,
            line_end=slide_index,
            locator={**locator, "source_image_id": source_image.image_id},
        )
    )


def _append_speaker_notes(
    slide: Any,
    heading_path: tuple[str, ...],
    slide_index: int,
    blocks: list[DocumentIRBlock],
) -> None:
    notes_frame = getattr(getattr(slide, "notes_slide", None), "notes_text_frame", None)
    notes = str(getattr(notes_frame, "text", "")).strip()
    if not notes:
        return
    blocks.append(
        DocumentIRBlock(
            block_id=_block_id(),
            ordinal=len(blocks),
            kind="paragraph",
            text=notes,
            heading_path=heading_path,
            line_start=slide_index,
            line_end=slide_index,
            locator={"slide": slide_index, "slide_index": slide_index, "source": "speaker_notes"},
        )
    )


def _markdown_table(values: list[list[str]]) -> str:
    width = max(len(row) for row in values)
    normalized = [row + [""] * (width - len(row)) for row in values]
    header = _markdown_row(normalized[0])
    separator = "| " + " | ".join("---" for _ in range(width)) + " |"
    return "\n".join([header, separator, *(_markdown_row(row) for row in normalized[1:])])


def _markdown_row(row: list[str]) -> str:
    return "| " + " | ".join(_markdown_cell(value) for value in row) + " |"


def _markdown_cell(value: str) -> str:
    return value.replace("|", "\\|").replace("\n", "<br>")


def _block_id() -> str:
    return uuid4().hex
